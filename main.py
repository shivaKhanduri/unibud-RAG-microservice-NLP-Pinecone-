import os
import io
import json
import time
import tempfile
import logging
import asyncio
from typing import List, Dict, Optional
from urllib.parse import quote

from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Response, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import openai  # <-- Using openai client for chat completions
import httpx
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import pytesseract
from google.cloud import storage
from pinecone import Pinecone, ServerlessSpec
from pydantic import BaseModel
import spacy

# ------------------------------------------------------------------------------
# Logging Configuration
# ------------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ------------------------------------------------------------------------------
# spaCy: Ensure en_core_web_sm is installed
# ------------------------------------------------------------------------------
try:
    nlp = spacy.load("en_core_web_sm")
except Exception as e:
    raise Exception("spaCy model 'en_core_web_sm' not found. Please run 'python -m spacy download en_core_web_sm'.")

# ------------------------------------------------------------------------------
# FastAPI Initialization
# ------------------------------------------------------------------------------
app = FastAPI()

@app.middleware("http")
async def add_process_time_header(request: Request, call_next):
    start_time = time.time()
    response: Response = await call_next(request)
    process_time = time.time() - start_time
    response.headers["X-Process-Time"] = str(process_time)
    logger.info(f"Request: {request.url.path} completed in {process_time:.4f} seconds")
    return response

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ------------------------------------------------------------------------------
# Environment Configuration
# ------------------------------------------------------------------------------
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")

# Configure openai to point to DeepSeek
# (DeepSeek's OpenAI-compatible base URL)
openai.api_key = DEEPSEEK_API_KEY
openai.api_base = "https://api.deepseek.com"

# Firebase configuration
firebase_creds_json = os.getenv("GOOGLE_APPLICATION_CREDENTIALS_JSON")
if firebase_creds_json:
    try:
        creds_data = json.loads(firebase_creds_json)
        with tempfile.NamedTemporaryFile(mode='w', delete=False) as f:
            json.dump(creds_data, f)
            os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = f.name
    except json.JSONDecodeError:
        raise Exception("Invalid JSON in GOOGLE_APPLICATION_CREDENTIALS_JSON")

# Initialize Firebase
storage_client = storage.Client()
bucket = storage_client.bucket("unibud-22153.appspot.com")

# Initialize Pinecone
pc = Pinecone(api_key=PINECONE_API_KEY)
INDEX_NAME = "unibud-index"
EMBEDDING_DIM = 1536  # Adjust if needed for your embedding dimension

# Create index if not exists
if INDEX_NAME not in pc.list_indexes().names():
    pc.create_index(
        name=INDEX_NAME,
        dimension=EMBEDDING_DIM,
        metric="cosine",
        spec=ServerlessSpec(cloud="aws", region="us-east-1")
    )

index = pc.Index(INDEX_NAME)

# ------------------------------------------------------------------------------
# Helper Functions
# ------------------------------------------------------------------------------
def upload_to_firebase(file_bytes: bytes, file_name: str, content_type: str) -> str:
    """Uploads file bytes to Firebase Storage."""
    blob = bucket.blob(f"study_resources/{file_name}")
    blob.upload_from_string(file_bytes, content_type=content_type)
    blob.make_public()
    return blob.public_url

def chunk_text(text: str, max_sentences: int = 5, overlap: int = 1) -> List[str]:
    """Split text into chunks using spaCy."""
    doc = nlp(text)
    sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
    chunks = []
    start = 0
    while start < len(sentences):
        end = start + max_sentences
        chunk = " ".join(sentences[start:end])
        if chunk:
            chunks.append(chunk)
        start += max(1, max_sentences - overlap)
    return chunks

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a PDF with OCR fallback."""
    doc = fitz.open(pdf_path)
    full_text = ""
    for page in doc:
        text = page.get_text().strip()
        if len(text) < 20:
            try:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_text = pytesseract.image_to_string(img)
                full_text += ocr_text + "\n"
            except Exception as e:
                logger.warning(f"OCR error on page {page.number}: {e}")
                full_text += "\n"
        else:
            full_text += text + "\n"
    return full_text

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extracts text from PPT/PPTX with OCR for images."""
    prs = Presentation(ppt_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                full_text += shape.text + "\n"
            # shape_type == 13 typically indicates a picture placeholder
            if hasattr(shape, "shape_type") and shape.shape_type == 13:
                try:
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    ocr_text = pytesseract.image_to_string(image)
                    full_text += ocr_text + "\n"
                except Exception as e:
                    logger.warning(f"Error processing image OCR in PPT: {e}")
    return full_text

def process_file(file_path: str) -> str:
    """Processes a file based on its extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    return ""

# ------------------------------------------------------------------------------
# Embeddings with httpx (if you still need them)
# ------------------------------------------------------------------------------
async def get_embeddings(texts: List[str]) -> List[List[float]]:
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(
                "https://api.deepseek.com/v1/embeddings",  # Added /v1
                headers={
                    "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "input": texts,
                    "model": "text-embedding-002"  # Verify model name
                },
                timeout=30
            )
            response.raise_for_status()
            data = response.json()
            
            return [item["embedding"] for item in data["data"]]
    except Exception as e:
        logger.error(f"DeepSeek embedding error: {e}")
        raise HTTPException(500, detail=f"Embedding error: {e}")

async def get_embedding(text: str) -> List[float]:
    """Convenience method for a single embedding."""
    embeddings = await get_embeddings([text])
    return embeddings[0]

# ------------------------------------------------------------------------------
# Chat Completion (OpenAI Python Client, pointing to DeepSeek)
# ------------------------------------------------------------------------------
def deepseek_chat_completion_sync(messages: List[Dict]) -> str:
    """
    Synchronous function that calls the DeepSeek chat model
    via the OpenAI Python client library. 
    """
    try:
        response = openai.ChatCompletion.create(
            model="deepseek-chat",
            messages=messages,
            temperature=0.2,
            max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        logger.error(f"DeepSeek chat error: {e}")
        raise HTTPException(500, detail=f"Chat completion error: {e}")

# ------------------------------------------------------------------------------
# Document Indexing
# ------------------------------------------------------------------------------
async def index_documents(file_infos: List[Dict]) -> None:
    """Processes and indexes documents with DeepSeek embeddings into Pinecone."""
    vectors = []
    for file_info in file_infos:
        text = process_file(file_info["local_path"])
        if not text.strip():
            continue

        chunks = chunk_text(text)
        if chunks:
            embeddings = await get_embeddings(chunks)
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                vector_id = f"{os.path.basename(file_info['local_path'])}_{i}"
                metadata = {
                    "source_file": os.path.basename(file_info["local_path"]),
                    "chunk_index": i,
                    "text": chunk,
                    "file_url": file_info["firebase_url"]
                }
                if "user_metadata" in file_info:
                    metadata.update(file_info["user_metadata"])
                vectors.append((vector_id, embedding, metadata))
                await asyncio.sleep(0.2)  # non-blocking sleep

    if vectors:
        index.upsert(vectors=vectors)
        logger.info(f"Upserted {len(vectors)} vectors into Pinecone.")

# ------------------------------------------------------------------------------
# API Endpoints
# ------------------------------------------------------------------------------
class QueryRequest(BaseModel):
    query: str

@app.post("/upload")
async def upload_files(
    files: List[UploadFile] = File(...),
    metadata: Optional[str] = Form(None)
):
    start_time = time.time()
    try:
        file_infos = []
        temp_paths = []

        # Parse any incoming metadata JSON
        user_metadata_list = []
        if metadata:
            try:
                user_metadata_list = json.loads(metadata)
                if not isinstance(user_metadata_list, list):
                    raise ValueError("Metadata should be a list of objects.")
            except Exception as e:
                raise HTTPException(400, detail=f"Invalid metadata JSON: {e}")

        # Process each uploaded file
        for idx, uploaded_file in enumerate(files):
            if not uploaded_file.filename.lower().endswith((".pdf", ".ppt", ".pptx")):
                raise HTTPException(400, detail="Invalid file type")

            content = await uploaded_file.read()
            firebase_url = upload_to_firebase(
                content,
                uploaded_file.filename,
                uploaded_file.content_type
            )

            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.filename)[1]) as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name

            file_info = {
                "local_path": temp_path,
                "firebase_url": firebase_url
            }
            if idx < len(user_metadata_list):
                file_info["user_metadata"] = user_metadata_list[idx]
            file_infos.append(file_info)
            temp_paths.append(temp_path)

        # Index documents into Pinecone
        await index_documents(file_infos)

        # Cleanup local temp files
        for path in temp_paths:
            os.unlink(path)

        processing_time = time.time() - start_time
        return JSONResponse(status_code=200, content={
            "message": f"Successfully processed {len(files)} file(s)",
            "processing_time": processing_time
        })

    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(500, detail=str(e))

@app.post("/ask")
async def ask_question(query: QueryRequest):
    """
    Asks a question:
      1. Gets an embedding of the query.
      2. Queries Pinecone for context.
      3. Calls the DeepSeek chat model (via openai client) to form an answer.
    """
    start_time = time.time()
    try:
        query_text = query.query.strip()
        if not query_text:
            raise HTTPException(400, detail="Empty query provided.")

        # 1. Get embedding from DeepSeek
        query_embedding = await get_embedding(query_text)

        # 2. Query Pinecone
        response = index.query(vector=query_embedding, top_k=10, include_metadata=True)

        # 3. Build context
        matches = [m for m in response.matches if m.metadata.get("text", "").strip()]
        if not matches:
            processing_time = time.time() - start_time
            return {
                "answer": "I couldn't find relevant information to answer your question.",
                "sources": [],
                "highlighted_context": "",
                "processing_time": processing_time,
                "token_usage": {}
            }

        context_chunks = [match.metadata["text"] for match in matches]
        context = "\n".join(context_chunks)
        highlighted_context = "\n".join([f"<mark>{chunk}</mark>" for chunk in context_chunks])

        # 4. Prepare messages for the chat model
        system_prompt = (
            "You are a knowledgeable teaching assistant. Use ONLY the context provided below to answer the question. "
            "Provide a detailed, step-by-step explanation. If the context does not contain the answer, explicitly state that no relevant information was found."
        )
        user_prompt = f"Context:\n{context}\n\nQuestion: {query_text}"
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        # 5. Because openai is synchronous, run in executor to avoid blocking
        loop = asyncio.get_event_loop()
        answer = await loop.run_in_executor(
            None,
            lambda: deepseek_chat_completion_sync(messages)
        )

        # 6. Gather sources (limit as needed)
        sources = [match.metadata.get("file_url", "") for match in matches]

        processing_time = time.time() - start_time
        return {
            "answer": answer,
            "sources": sources[:2],
            "context": context,
            "highlighted_context": highlighted_context,
            "processing_time": processing_time,
            "token_usage": {}
        }

    except Exception as e:
        logger.error(f"Ask error: {e}")
        raise HTTPException(500, detail=f"Query processing error: {e}")

@app.get("/health")
def health_check():
    return {"status": "healthy"}

@app.get("/")
def read_root():
    return {"message": "UniBud API - DeepSeek Edition"}
